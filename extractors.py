# extractors.py
import re, ast
from pathlib import Path
import PyPDF2, docx
from pptx import Presentation

ALLOWED_DOCS = {".pdf", ".docx", ".txt", ".md", ".pptx"}
ALLOWED_SCRIPTS = {".py", ".js", ".java", ".cpp", ".c", ".r", ".ts", ".html", ".css"}

def read_textlike(path: str) -> str:
    ext = Path(path).suffix.lower()
    if ext in [".txt", ".md"]:
        return Path(path).read_text(encoding="utf-8", errors="ignore")
    if ext == ".docx":
        d = docx.Document(path); return "\n".join(p.text for p in d.paragraphs)
    if ext == ".pdf":
        with open(path, "rb") as f:
            r = PyPDF2.PdfReader(f)
            return "\n".join((p.extract_text() or "") for p in r.pages)
    if ext == ".pptx":
        pres = Presentation(path)
        bits = []
        for i, slide in enumerate(pres.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    bits.append(f"[Slide {i+1}] {shape.text}")
        return "\n".join(bits)
    raise ValueError(f"Unsupported doc type: {ext}")

def split_code_and_comments(path: str):
    ext = Path(path).suffix.lower()
    raw = Path(path).read_text(encoding="utf-8", errors="ignore")

    # Python: hash comments + docstrings via AST
    if ext == ".py":
        comments, code_nc_lines = [], []
        for line in raw.splitlines():
            if re.match(r'^\s*#', line):
                comments.append(line.lstrip()); code_nc_lines.append("")
            else:
                code_nc_lines.append(line)
        try:
            t = ast.parse(raw); docstrings = []
            if ast.get_docstring(t): docstrings.append(ast.get_docstring(t))
            for node in ast.walk(t):
                if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef, ast.ClassDef)):
                    ds = ast.get_docstring(node); 
                    if ds: docstrings.append(ds)
            for ds in docstrings: comments.append(f'""" {ds} """')
        except: pass
        return {"language":"python","raw":raw,
                "code_no_comments":"\n".join(code_nc_lines),
                "comments_only":"\n".join(comments)}

    # C/C++/Java/JS/TS/CSS: // and /* ... */
    if ext in {".js",".ts",".java",".cpp",".c",".css"}:
        line_comment = re.compile(r'^\s*//(.*)$', re.MULTILINE)
        block_comment = re.compile(r'/\*.*?\*/', re.DOTALL)
        comments = []
        comments.extend([m.group(0) for m in line_comment.finditer(raw)])
        comments.extend(block_comment.findall(raw))
        code_nc = block_comment.sub("", line_comment.sub("", raw))
        lang = {".js":"javascript",".ts":"typescript",".java":"java",".cpp":"cpp",".c":"c",".css":"css"}[ext]
        return {"language":lang,"raw":raw,"code_no_comments":code_nc,"comments_only":"\n".join(comments)}

    # HTML: <!-- ... -->
    if ext == ".html":
        html_comments = re.findall(r'<!--(.*?)-->', raw, re.DOTALL)
        code_nc = re.sub(r'<!--.*?-->', "", raw, flags=re.DOTALL)
        return {"language":"html","raw":raw,"code_no_comments":code_nc,"comments_only":"\n".join(["<!--"+c+"-->" for c in html_comments])}

    # R: '#'
    if ext == ".r":
        comments, code_nc_lines = [], []
        for line in raw.splitlines():
            if re.match(r'^\s*#', line):
                comments.append(line.lstrip()); code_nc_lines.append("")
            else:
                code_nc_lines.append(line)
        return {"language":"r","raw":raw,"code_no_comments":"\n".join(code_nc_lines),"comments_only":"\n".join(comments)}

    raise ValueError(f"Unsupported script type: {ext}")
